"""
Presentation Service — Professional PPTX generation with AI-generated visuals.
Uses python-pptx for slides, Pollinations.ai for image generation,
and SearXNG deep research for content enrichment.

Enhanced with MINI/MIDI/MAXI mode analysis, multi-theme support,
and advanced slide types (quote, data highlight, two-column).
"""

from __future__ import annotations

import asyncio
import io
import json
import re
import tempfile
import urllib.parse
from pathlib import Path
from typing import Any, Literal

import httpx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Presentation Modes ───────────────────────────────────────────

PresentationMode = Literal["mini", "midi", "maxi"]

MODE_CONFIG: dict[PresentationMode, dict[str, Any]] = {
    "mini": {
        "slide_range": (5, 7),
        "default_slides": 6,
        "research_queries": 3,
        "research_depth": 2,       # top N pages to fetch full content (reduced for faster response)
        "label": "MINI",
        "emoji": "⚡",
        "description_tr": "Özet sunum — ana noktalar ve temel veriler",
        "description_en": "Executive summary — key points and core data",
    },
    "midi": {
        "slide_range": (10, 15),
        "default_slides": 12,
        "research_queries": 5,
        "research_depth": 3,
        "label": "MIDI",
        "emoji": "📊",
        "description_tr": "Standart sunum — detaylı analiz, örnekler ve veriler",
        "description_en": "Standard presentation — detailed analysis, examples and data",
    },
    "maxi": {
        "slide_range": (20, 30),
        "default_slides": 25,
        "research_queries": 10,
        "research_depth": 6,
        "label": "MAXI",
        "emoji": "🎯",
        "description_tr": "Kapsamlı sunum — derinlemesine araştırma, vaka çalışmaları, karşılaştırmalar",
        "description_en": "Comprehensive presentation — deep research, case studies, comparisons",
    },
}


# ── Theme System ─────────────────────────────────────────────────

THEMES: dict[str, dict[str, RGBColor]] = {
    "corporate": {
        "primary": RGBColor(0x0F, 0x34, 0x60),
        "secondary": RGBColor(0x53, 0x34, 0x83),
        "accent": RGBColor(0xE9, 0x4D, 0x60),
        "bg_dark": RGBColor(0x16, 0x21, 0x3E),
        "bg_light": RGBColor(0xF8, 0xF9, 0xFA),
        "text_dark": RGBColor(0x1A, 0x1A, 0x2E),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x99, 0x99, 0x99),
        "green": RGBColor(0x10, 0xB9, 0x81),
        "highlight": RGBColor(0x3B, 0x82, 0xF6),
    },
    "modern_dark": {
        "primary": RGBColor(0x6C, 0x5C, 0xE7),
        "secondary": RGBColor(0xA2, 0x9B, 0xFE),
        "accent": RGBColor(0xFD, 0x79, 0xA8),
        "bg_dark": RGBColor(0x0D, 0x0D, 0x0D),
        "bg_light": RGBColor(0x1A, 0x1A, 0x2E),
        "text_dark": RGBColor(0xE0, 0xE0, 0xE0),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x88, 0x88, 0x99),
        "green": RGBColor(0x00, 0xCE, 0xC9),
        "highlight": RGBColor(0xFD, 0xCB, 0x6E),
    },
    "nature": {
        "primary": RGBColor(0x00, 0x69, 0x5C),
        "secondary": RGBColor(0x00, 0x89, 0x7B),
        "accent": RGBColor(0xFF, 0x8F, 0x00),
        "bg_dark": RGBColor(0x00, 0x2B, 0x25),
        "bg_light": RGBColor(0xF1, 0xF8, 0xE9),
        "text_dark": RGBColor(0x1B, 0x2A, 0x1B),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x7C, 0x9A, 0x7C),
        "green": RGBColor(0x4C, 0xAF, 0x50),
        "highlight": RGBColor(0xFF, 0xC1, 0x07),
    },
}

# Default theme (backward compatible)
COLORS = THEMES["corporate"]

# ── Image Generation (Pollinations.ai — authenticated with API key) ──

POLLINATIONS_BASE = "https://gen.pollinations.ai/image"
POLLINATIONS_MODEL = "zimage"
POLLINATIONS_API_KEY = "sk_IPFPaqTQv7Jpw4SDTTOKtF8qwN4SFZgZ"


# ── Deep Research for Presentations ──────────────────────────────

async def deep_research_for_presentation(
    topic: str,
    language: str = "tr",
    max_queries: int = 5,
    mode: PresentationMode = "midi",
) -> dict[str, Any]:
    """
    Perform multi-query deep research on a topic using SearXNG.
    Research depth scales with mode: MINI=light, MIDI=standard, MAXI=deep.
    Returns structured research data: key_facts, statistics, sources, subtopics.
    """
    from tools.search import web_search

    cfg = MODE_CONFIG[mode]
    effective_queries = cfg["research_queries"]
    depth = cfg["research_depth"]

    # Generate diverse search queries for comprehensive coverage
    queries = _generate_research_queries(topic, language, mode)[:effective_queries]

    # Run all searches in parallel (fewer results per query for faster MINI/MIDI)
    results_per_query = 4 if mode == "mini" else 5 if mode == "midi" else 10
    tasks = [web_search(q, max_results=results_per_query) for q in queries]
    all_results = await asyncio.gather(*tasks, return_exceptions=True)

    # Aggregate and deduplicate
    seen_urls: set[str] = set()
    facts: list[str] = []
    sources: list[dict[str, str]] = []
    all_snippets: list[str] = []

    for result_set in all_results:
        if isinstance(result_set, Exception):
            continue
        for item in result_set:
            url = item.get("url", "")
            if url in seen_urls or not url:
                continue
            seen_urls.add(url)
            snippet = item.get("snippet", "").strip()
            if snippet and len(snippet) > 30:
                facts.append(snippet)
                all_snippets.append(snippet)
                sources.append({
                    "title": item.get("title", ""),
                    "url": url,
                })

    # Try to fetch detailed content from top sources (depth scales with mode)
    detailed_content: list[str] = []
    top_urls = [s["url"] for s in sources[:depth]]
    max_content_chars = 2000 if mode == "mini" else 3000 if mode == "midi" else 5000
    fetch_tasks = [_fetch_page_content(url, max_chars=max_content_chars) for url in top_urls]
    fetch_results = await asyncio.gather(*fetch_tasks, return_exceptions=True)
    for content in fetch_results:
        if isinstance(content, str) and len(content) > 100:
            detailed_content.append(content[:max_content_chars])

    return {
        "queries_used": queries,
        "facts": facts[:30 if mode == "mini" else 50 if mode == "midi" else 80],
        "sources": sources[:15 if mode == "mini" else 25 if mode == "midi" else 40],
        "detailed_content": detailed_content,
        "total_sources": len(sources),
        "mode": mode,
    }


def _generate_research_queries(topic: str, language: str = "tr",
                               mode: PresentationMode = "midi") -> list[str]:
    """Generate diverse search queries for comprehensive topic coverage.
    More queries for MAXI mode, fewer for MINI."""
    base = topic.strip()

    if language == "tr":
        queries = [
            base,
            f"{base} istatistik veri 2024 2025",
            f"{base} avantajları dezavantajları",
            f"{base} örnekler kullanım alanları",
            f"{base} gelecek trendler",
        ]
        if mode in ("midi", "maxi"):
            queries.extend([
                f"{base} pazar büyüklüğü rapor",
                f"{base} karşılaştırma alternatifleri",
                f"{base} başarı hikayeleri vaka çalışması",
            ])
        if mode == "maxi":
            queries.extend([
                f"{base} teknik detaylar mimari",
                f"{base} zorluklar riskler çözümler",
                f"{base} akademik araştırma makale",
                f"{base} dünya örnekleri lider şirketler",
            ])
    else:
        queries = [
            base,
            f"{base} statistics data 2024 2025",
            f"{base} advantages disadvantages",
            f"{base} examples use cases",
            f"{base} future trends",
        ]
        if mode in ("midi", "maxi"):
            queries.extend([
                f"{base} market size report",
                f"{base} comparison alternatives",
                f"{base} success stories case study",
            ])
        if mode == "maxi":
            queries.extend([
                f"{base} technical details architecture",
                f"{base} challenges risks solutions",
                f"{base} academic research paper",
                f"{base} world examples leading companies",
            ])

    return queries


async def _fetch_page_content(url: str, max_chars: int = 3000) -> str:
    """Fetch and extract text content from a URL."""
    try:
        async with httpx.AsyncClient(timeout=8.0, follow_redirects=True) as client:
            resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0"})
            if resp.status_code != 200:
                return ""
            # Simple HTML text extraction
            text = re.sub(r"<script[^>]*>.*?</script>", "", resp.text, flags=re.DOTALL)
            text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
            text = re.sub(r"<[^>]+>", " ", text)
            text = re.sub(r"\s+", " ", text).strip()
            return text[:max_chars]
    except Exception:
        return ""


def format_research_for_prompt(research: dict[str, Any]) -> str:
    """Format research results into a structured context block for the LLM."""
    parts = ["<deep_research>"]

    if research.get("facts"):
        parts.append("\n## Araştırma Bulguları:")
        for i, fact in enumerate(research["facts"][:20], 1):
            parts.append(f"  {i}. {fact}")

    if research.get("detailed_content"):
        parts.append("\n## Detaylı İçerik:")
        for i, content in enumerate(research["detailed_content"], 1):
            parts.append(f"  --- Kaynak {i} ---")
            parts.append(f"  {content[:1500]}")

    if research.get("sources"):
        parts.append("\n## Kaynaklar:")
        for s in research["sources"][:10]:
            parts.append(f"  - {s['title']}: {s['url']}")

    parts.append("</deep_research>")
    return "\n".join(parts)


def _get_pollinations_token() -> str:
    """Load Pollinations API key — env override veya hardcoded key."""
    import os
    return (
        os.environ.get("POLLINATIONS_API_KEY")
        or os.environ.get("Pollinations_api_key")
        or POLLINATIONS_API_KEY
    )


def _build_image_url(prompt: str, width: int, height: int) -> str:
    """Build Pollinations gen API URL with zimage model."""
    encoded = urllib.parse.quote(prompt)
    params = {
        "model": POLLINATIONS_MODEL,
        "width": str(width),
        "height": str(height),
        "nologo": "true",
        "enhance": "true",
    }
    qs = urllib.parse.urlencode(params)
    return f"{POLLINATIONS_BASE}/{encoded}?{qs}"


async def generate_image(prompt: str, width: int = 1280, height: int = 720) -> bytes | None:
    """
    Generate an image using Pollinations gen API (zimage model) with retry and fallback.
    Returns image bytes or None on failure.
    """
    token = _get_pollinations_token()
    headers = {"Authorization": f"Bearer {token}"} if token else {}

    for attempt in range(2):
        url = _build_image_url(prompt, width, height)
        try:
            async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
                resp = await client.get(url, headers=headers)
                if resp.status_code == 200 and len(resp.content) > 1000:
                    return resp.content
                print(f"[PresentationService] Image attempt {attempt+1}: status={resp.status_code}, size={len(resp.content)}")
        except Exception as e:
            print(f"[PresentationService] Image attempt {attempt+1} failed: {e}")

        # Simplify prompt for retry
        if attempt == 0:
            prompt = prompt[:100] + ", professional, clean design"

    # Fallback: simpler prompt ile tekrar dene
    try:
        simple = re.sub(r"[^\w\s]", "", prompt.split(",")[0])[:50]
        fallback_url = _build_image_url(f"professional photo {simple}", width, height)
        async with httpx.AsyncClient(timeout=45.0, follow_redirects=True) as client:
            resp = await client.get(fallback_url, headers=headers)
            if resp.status_code == 200 and len(resp.content) > 1000:
                return resp.content
    except Exception:
        pass

    return None


def _generate_image_sync(prompt: str, width: int = 1280, height: int = 720) -> bytes | None:
    """Sync version of image generation with retry."""
    token = _get_pollinations_token()
    headers = {"Authorization": f"Bearer {token}"} if token else {}
    for attempt in range(2):
        url = _build_image_url(prompt, width, height)
        try:
            with httpx.Client(timeout=60.0, follow_redirects=True) as client:
                resp = client.get(url, headers=headers)
                if resp.status_code == 200 and len(resp.content) > 1000:
                    return resp.content
                print(f"[PresentationService] Sync image attempt {attempt+1}: status={resp.status_code}")
        except Exception as e:
            print(f"[PresentationService] Sync image attempt {attempt+1} failed: {e}")
        if attempt == 0:
            prompt = prompt[:100] + ", professional, clean design"
    return None


# ── Slide Builders ───────────────────────────────────────────────

def _set_slide_bg(slide, color: RGBColor):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text: str,
                  font_size: int = 18, color: RGBColor = None,
                  bold: bool = False, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or COLORS["text_dark"]
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _build_title_slide(prs: Presentation, title: str, subtitle: str = "",
                       colors: dict | None = None):
    """Create a cinematic title slide with geometric accents."""
    c = colors or COLORS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, c["bg_dark"])

    # ── Left vertical accent bar (bold color strip) ──
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.35), Inches(7.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = c["accent"]
    bar.line.fill.background()

    # ── Top-right diagonal decorative triangle ──
    tri = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(10.33), Inches(0), Inches(3), Inches(2.2),
    )
    tri.fill.solid()
    tri.fill.fore_color.rgb = c["primary"]
    tri.line.fill.background()
    tri.rotation = 0

    # ── Bottom decorative band ──
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.6), Inches(13.33), Inches(0.9),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = c["primary"]
    band.line.fill.background()

    # ── Small accent square (geometric detail) ──
    sq = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0), Inches(2.0), Inches(0.5), Inches(0.5),
    )
    sq.fill.solid()
    sq.fill.fore_color.rgb = c["accent"]
    sq.line.fill.background()

    # ── Title ──
    _add_text_box(
        slide, Inches(1.2), Inches(2.7), Inches(10.5), Inches(1.8),
        title, font_size=40, color=c["text_light"],
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    # ── Accent underline below title ──
    uline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2), Inches(4.5), Inches(3), Inches(0.06),
    )
    uline.fill.solid()
    uline.fill.fore_color.rgb = c["accent"]
    uline.line.fill.background()

    # ── Subtitle ──
    if subtitle:
        _add_text_box(
            slide, Inches(1.2), Inches(4.8), Inches(10.5), Inches(1),
            subtitle, font_size=18, color=c["text_muted"],
            alignment=PP_ALIGN.LEFT,
        )

    # ── Footer on bottom band ──
    from datetime import datetime
    date_str = datetime.now().strftime("%d %B %Y")
    _add_text_box(
        slide, Inches(1.2), Inches(6.7), Inches(11), Inches(0.5),
        f"Multi-Agent Ops Center  •  {date_str}",
        font_size=11, color=c["text_muted"], alignment=PP_ALIGN.LEFT,
    )


def _build_content_slide(prs: Presentation, title: str, bullets: list[str],
                         slide_num: int, total_slides: int,
                         image_bytes: bytes | None = None,
                         colors: dict | None = None):
    """Create a polished content slide with accent shapes and visual hierarchy."""
    c = colors or COLORS
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Alternating background: even slides get light, odd get subtle tint ──
    if slide_num % 2 == 0:
        _set_slide_bg(slide, c["bg_light"])
        txt_color = c["text_dark"]
    else:
        _set_slide_bg(slide, c["bg_light"])
        txt_color = c["text_dark"]

    # ── Left accent strip ──
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.12), Inches(7.5),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = c["primary"]
    strip.line.fill.background()

    # ── Title background block ──
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.12), Inches(0), Inches(13.21), Inches(1.15),
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = c["primary"]
    title_bg.line.fill.background()

    # ── Title text (white on primary) ──
    _add_text_box(
        slide, Inches(0.7), Inches(0.15), Inches(11.5), Inches(0.8),
        title, font_size=24, color=c["text_light"], bold=True,
    )

    # ── Accent dot next to title ──
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.35), Inches(0.35), Inches(0.18), Inches(0.18),
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = c["accent"]
    dot.line.fill.background()

    # ── Content area ──
    if image_bytes:
        content_width = Inches(6.5)
        content_left = Inches(0.7)
        try:
            img_stream = io.BytesIO(image_bytes)
            # Image with rounded-look container
            img_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(7.3), Inches(1.5), Inches(5.6), Inches(4.3),
            )
            img_bg.fill.solid()
            img_bg.fill.fore_color.rgb = c["bg_dark"]
            img_bg.line.fill.background()
            slide.shapes.add_picture(
                img_stream, Inches(7.45), Inches(1.65),
                width=Inches(5.3), height=Inches(4),
            )
        except Exception:
            content_width = Inches(12)
    else:
        content_width = Inches(12)
        content_left = Inches(0.7)

    # ── Bullet points with custom markers ──
    bullet_markers = ["▸", "▹", "◆", "◇", "●", "○"]
    txBox = slide.shapes.add_textbox(
        content_left, Inches(1.5), content_width, Inches(5.2),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        marker = bullet_markers[i % len(bullet_markers)]
        p.text = f"{marker}  {bullet}"
        p.font.size = Pt(14)
        p.font.color.rgb = txt_color
        p.space_after = Pt(10)
        p.space_before = Pt(2)
        p.level = 0

    # ── Bottom accent line ──
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7), Inches(7.05), Inches(12), Inches(0.03),
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = c["accent"]
    bottom_line.line.fill.background()

    # ── Slide number ──
    _add_text_box(
        slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.35),
        f"{slide_num}/{total_slides}",
        font_size=9, color=c["text_muted"], alignment=PP_ALIGN.RIGHT,
    )


def _build_section_slide(prs: Presentation, title: str, slide_num: int, total_slides: int,
                         colors: dict | None = None):
    """Create a bold section divider slide with geometric design."""
    c = colors or COLORS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, c["bg_dark"])

    # ── Large diagonal accent block (top-left) ──
    diag = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(-1), Inches(-0.5), Inches(7), Inches(8.5),
    )
    diag.fill.solid()
    diag.fill.fore_color.rgb = c["primary"]
    diag.line.fill.background()

    # ── Accent horizontal line (center) ──
    hline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.5), Inches(3.2), Inches(4.5), Inches(0.06),
    )
    hline.fill.solid()
    hline.fill.fore_color.rgb = c["accent"]
    hline.line.fill.background()

    # ── Section title (large, centered) ──
    _add_text_box(
        slide, Inches(2), Inches(3.5), Inches(9.5), Inches(1.5),
        title, font_size=36, color=c["text_light"],
        bold=True, alignment=PP_ALIGN.CENTER,
    )

    # ── Accent underline below title ──
    uline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(5.1), Inches(2.5), Inches(0.06),
    )
    uline.fill.solid()
    uline.fill.fore_color.rgb = c["accent"]
    uline.line.fill.background()

    # ── Small decorative circle (bottom-right) ──
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11.5), Inches(6.0), Inches(0.6), Inches(0.6),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = c["accent"]
    circ.line.fill.background()

    # ── Slide number ──
    _add_text_box(
        slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.35),
        f"{slide_num}/{total_slides}",
        font_size=9, color=c["text_muted"], alignment=PP_ALIGN.RIGHT,
    )


def _build_closing_slide(prs: Presentation, title: str = "Teşekkürler",
                         colors: dict | None = None):
    """Create a cinematic closing slide mirroring the title slide style."""
    c = colors or COLORS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, c["bg_dark"])

    # ── Top band ──
    band_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9),
    )
    band_top.fill.solid()
    band_top.fill.fore_color.rgb = c["primary"]
    band_top.line.fill.background()

    # ── Bottom band ──
    band_bot = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.6), Inches(13.33), Inches(0.9),
    )
    band_bot.fill.solid()
    band_bot.fill.fore_color.rgb = c["primary"]
    band_bot.line.fill.background()

    # ── Accent circle (center decorative) ──
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6.17), Inches(1.8), Inches(1), Inches(1),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = c["accent"]
    circ.line.fill.background()

    # ── Title ──
    _add_text_box(
        slide, Inches(1), Inches(3.0), Inches(11.33), Inches(1.5),
        title, font_size=44, color=c["text_light"],
        bold=True, alignment=PP_ALIGN.CENTER,
    )

    # ── Accent underline ──
    uline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(4.6), Inches(2.5), Inches(0.06),
    )
    uline.fill.solid()
    uline.fill.fore_color.rgb = c["accent"]
    uline.line.fill.background()

    # ── Subtitle ──
    _add_text_box(
        slide, Inches(1), Inches(5.0), Inches(11.33), Inches(0.5),
        "Multi-Agent Ops Center ile oluşturuldu",
        font_size=13, color=c["text_muted"], alignment=PP_ALIGN.CENTER,
    )


# ── New Slide Types (MAXI mode) ─────────────────────────────────

def _build_quote_slide(prs: Presentation, title: str,
                       quote_text: str, quote_author: str,
                       bullets: list[str] | None = None,
                       slide_num: int = 0, total_slides: int = 0,
                       image_bytes: bytes | None = None,
                       colors: dict | None = None):
    """Create a dramatic quote slide with oversized quotation mark and accent styling."""
    c = colors or COLORS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, c["bg_dark"])

    # ── Left accent strip ──
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.12), Inches(7.5),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = c["accent"]
    strip.line.fill.background()

    # ── Title ──
    _add_text_box(
        slide, Inches(0.7), Inches(0.3), Inches(12), Inches(0.7),
        title, font_size=22, color=c["accent"], bold=True,
    )

    # ── Giant quotation mark (decorative) ──
    _add_text_box(
        slide, Inches(1.0), Inches(1.0), Inches(2), Inches(2),
        "\u201C", font_size=96, color=c["primary"],
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    # ── Quote text ──
    _add_text_box(
        slide, Inches(2.0), Inches(1.8), Inches(9.5), Inches(2.2),
        quote_text,
        font_size=22, color=c["text_light"], alignment=PP_ALIGN.LEFT,
    )

    # ── Accent line before author ──
    aline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0), Inches(4.1), Inches(1.5), Inches(0.04),
    )
    aline.fill.solid()
    aline.fill.fore_color.rgb = c["accent"]
    aline.line.fill.background()

    # ── Author ──
    _add_text_box(
        slide, Inches(2.0), Inches(4.3), Inches(9.5), Inches(0.5),
        f"— {quote_author}",
        font_size=14, color=c["text_muted"], alignment=PP_ALIGN.LEFT,
    )

    # ── Optional bullets below quote ──
    if bullets:
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(5.2), Inches(11), Inches(2),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets[:4]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▹  {bullet}"
            p.font.size = Pt(13)
            p.font.color.rgb = c["text_muted"]
            p.space_after = Pt(6)

    # ── Slide number ──
    _add_text_box(
        slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.35),
        f"{slide_num}/{total_slides}",
        font_size=9, color=c["text_muted"], alignment=PP_ALIGN.RIGHT,
    )


def _build_data_slide(prs: Presentation, title: str,
                      data_highlights: list[dict[str, str]],
                      bullets: list[str] | None = None,
                      slide_num: int = 0, total_slides: int = 0,
                      image_bytes: bytes | None = None,
                      colors: dict | None = None):
    """Create a data-focused slide with bold metric cards and accent styling."""
    c = colors or COLORS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, c["bg_light"])

    # ── Left accent strip ──
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.12), Inches(7.5),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = c["accent"]
    strip.line.fill.background()

    # ── Title background block ──
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.12), Inches(0), Inches(13.21), Inches(1.15),
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = c["bg_dark"]
    title_bg.line.fill.background()

    # ── Title ──
    _add_text_box(
        slide, Inches(0.7), Inches(0.15), Inches(11.5), Inches(0.8),
        title, font_size=24, color=c["text_light"], bold=True,
    )

    # ── Data highlight cards ──
    card_count = min(len(data_highlights), 3)
    card_width = 3.5
    total_card_width = card_count * card_width + (card_count - 1) * 0.6
    start_x = (13.33 - total_card_width) / 2

    card_colors = [c["primary"], c["accent"], c.get("green", c["primary"])]

    for idx, dh in enumerate(data_highlights[:3]):
        x = start_x + idx * (card_width + 0.6)
        card_color = card_colors[idx % len(card_colors)]

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.5), Inches(card_width), Inches(2.4),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_color
        card.line.fill.background()

        # Top accent line on card
        card_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.5), Inches(1.65), Inches(card_width - 1.0), Inches(0.04),
        )
        card_line.fill.solid()
        card_line.fill.fore_color.rgb = c["text_light"]
        card_line.line.fill.background()

        # Metric value (large)
        _add_text_box(
            slide, Inches(x + 0.2), Inches(1.85), Inches(card_width - 0.4), Inches(1),
            dh.get("value", ""),
            font_size=36, color=c["text_light"], bold=True, alignment=PP_ALIGN.CENTER,
        )
        # Metric name
        _add_text_box(
            slide, Inches(x + 0.2), Inches(2.85), Inches(card_width - 0.4), Inches(0.5),
            dh.get("metric", ""),
            font_size=13, color=c["text_light"], alignment=PP_ALIGN.CENTER,
        )
        # Context
        if dh.get("context"):
            _add_text_box(
                slide, Inches(x + 0.2), Inches(3.35), Inches(card_width - 0.4), Inches(0.4),
                dh["context"],
                font_size=10, color=c["text_muted"], alignment=PP_ALIGN.CENTER,
            )

    # ── Bullets below cards ──
    if bullets:
        txBox = slide.shapes.add_textbox(
            Inches(0.7), Inches(4.3), Inches(12), Inches(2.8),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets[:5]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸  {bullet}"
            p.font.size = Pt(13)
            p.font.color.rgb = c["text_dark"]
            p.space_after = Pt(7)

    # ── Image on right if available and space permits ──
    if image_bytes and not bullets:
        try:
            img_stream = io.BytesIO(image_bytes)
            slide.shapes.add_picture(
                img_stream, Inches(8.5), Inches(4.3),
                width=Inches(4.3), height=Inches(3),
            )
        except Exception:
            pass

    # ── Bottom accent line ──
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7), Inches(7.05), Inches(12), Inches(0.03),
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = c["accent"]
    bottom_line.line.fill.background()

    # ── Slide number ──
    _add_text_box(
        slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.35),
        f"{slide_num}/{total_slides}",
        font_size=9, color=c["text_muted"], alignment=PP_ALIGN.RIGHT,
    )


# ── Slide Content Parser ─────────────────────────────────────────

def parse_slide_content(raw_content: str) -> list[dict[str, Any]]:
    """
    Parse agent-generated content into structured slide data.
    Supports: SLIDE, SECTION, IMAGE, QUOTE, DATA markers.
    Also handles markdown headers (## Slide 1: Title) and other LLM formats.
    """
    slides = []
    current_slide = None

    for line in raw_content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue

        # New slide marker — strict format
        slide_match = re.match(r"^(?:SLIDE|SLAYT)\s*(\d+)\s*[:：]\s*(.+)", stripped, re.IGNORECASE)
        if slide_match:
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                "num": int(slide_match.group(1)),
                "title": slide_match.group(2).strip(),
                "bullets": [],
                "image_prompt": None,
                "is_section": False,
                "quote": None,
                "data_highlights": [],
            }
            continue

        # Markdown header slide format: ## Slide 1: Title or ### 1. Title
        md_slide_match = re.match(r"^#{1,4}\s*(?:Slide|Slayt|Slayd)?\s*(\d+)\s*[:.)\-—]\s*(.+)", stripped, re.IGNORECASE)
        if md_slide_match:
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                "num": int(md_slide_match.group(1)),
                "title": md_slide_match.group(2).strip().rstrip("#").strip(),
                "bullets": [],
                "image_prompt": None,
                "is_section": False,
                "quote": None,
                "data_highlights": [],
            }
            continue

        # Markdown header as section: ## Section Title (no number)
        md_section_match = re.match(r"^#{2,3}\s+(?!Slide|Slayt|Slayd|\d)(.{5,80})$", stripped, re.IGNORECASE)
        if md_section_match and current_slide:
            # Only treat as section if we already have slides
            if slides:
                slides.append(current_slide)
                current_slide = {
                    "num": len(slides) + 1,
                    "title": md_section_match.group(1).strip().rstrip("#").strip(),
                    "bullets": [],
                    "image_prompt": None,
                    "is_section": True,
                    "quote": None,
                    "data_highlights": [],
                }
                continue

        # Section marker
        section_match = re.match(r"^(?:SECTION|BÖLÜM)\s*[:：]\s*(.+)", stripped, re.IGNORECASE)
        if section_match:
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                "num": len(slides) + 1,
                "title": section_match.group(1).strip(),
                "bullets": [],
                "image_prompt": None,
                "is_section": True,
                "quote": None,
                "data_highlights": [],
            }
            continue

        # Image prompt
        img_match = re.match(r"^(?:IMAGE|GÖRSEL|RESIM|Image)\s*[:：]\s*(.+)", stripped, re.IGNORECASE)
        if img_match and current_slide:
            current_slide["image_prompt"] = img_match.group(1).strip()
            continue

        # Quote marker (MAXI mode)
        quote_match = re.match(r'^(?:QUOTE|ALINTI)\s*[:：]\s*["\u201c](.+?)["\u201d]\s*[-—]\s*(.+)', stripped, re.IGNORECASE)
        if quote_match and current_slide:
            current_slide["quote"] = {
                "text": quote_match.group(1).strip(),
                "author": quote_match.group(2).strip(),
            }
            continue

        # Data highlight marker (MAXI mode)
        data_match = re.match(r"^(?:DATA|VERİ)\s*[:：]\s*(.+)", stripped, re.IGNORECASE)
        if data_match and current_slide:
            parts = [p.strip() for p in data_match.group(1).split("|")]
            if len(parts) >= 2:
                current_slide["data_highlights"].append({
                    "metric": parts[0],
                    "value": parts[1],
                    "context": parts[2] if len(parts) > 2 else "",
                })
            continue

        # Bullet point
        if current_slide and stripped.startswith(("-", "•", "*")):
            bullet_text = stripped.lstrip("-•* ").strip()
            if bullet_text:
                current_slide["bullets"].append(bullet_text)
            continue

        # Numbered item
        num_match = re.match(r"^\d+[.)]\s*(.+)", stripped)
        if num_match and current_slide:
            current_slide["bullets"].append(num_match.group(1).strip())
            continue

        # Bold markdown as bullet: **text**
        bold_match = re.match(r"^\*\*(.+?)\*\*(.*)$", stripped)
        if bold_match and current_slide:
            text = bold_match.group(1).strip()
            rest = bold_match.group(2).strip().lstrip(":- ").strip()
            if rest:
                text = f"{text}: {rest}"
            current_slide["bullets"].append(text)
            continue

        # Plain text as bullet (only if reasonably sized)
        if current_slide and 5 < len(stripped) < 200:
            current_slide["bullets"].append(stripped)

    if current_slide:
        slides.append(current_slide)

    # Final aggressive fallback: if no slides were parsed but content exists,
    # try splitting by double newlines and creating slides from paragraphs
    if not slides and raw_content and len(raw_content.strip()) > 50:
        print(f"[parse_slide_content] No slides parsed from {len(raw_content)} chars, trying paragraph split...")
        paragraphs = [p.strip() for p in re.split(r"\n\s*\n", raw_content) if p.strip() and len(p.strip()) > 10]
        for i, para in enumerate(paragraphs[:30], 1):
            lines = [l.strip() for l in para.split("\n") if l.strip()]
            if not lines:
                continue
            # First line as title, rest as bullets
            title_line = lines[0][:80].lstrip("#-•*0123456789.) ").strip()
            if not title_line or len(title_line) < 3:
                title_line = f"Slayt {i}"
            bullets = []
            for l in lines[1:6]:
                clean = l.lstrip("#-•*0123456789.) ").strip()
                if clean and len(clean) > 3:
                    bullets.append(clean[:150])
            if not bullets and len(lines) > 0:
                bullets = [title_line[:100]]
                title_line = f"Slayt {i}"
            slides.append({
                "num": i,
                "title": title_line,
                "bullets": bullets,
                "image_prompt": None,
                "is_section": False,
                "quote": None,
                "data_highlights": [],
            })
        print(f"[parse_slide_content] Paragraph fallback produced {len(slides)} slides")

    return slides


# ── Main PPTX Generator ─────────────────────────────────────────

async def generate_presentation(
    slides_data: list[dict[str, Any]],
    title: str,
    subtitle: str = "",
    with_images: bool = True,
    theme: str = "corporate",
) -> bytes:
    """
    Generate a professional PPTX presentation.
    
    Args:
        slides_data: List of slide dicts with title, bullets, image_prompt, quote, data_highlights
        title: Presentation title
        subtitle: Presentation subtitle
        with_images: Whether to generate images via Pollinations.ai
        theme: Theme name (corporate, modern_dark, nature)
    
    Returns:
        PPTX file as bytes
    """
    colors = THEMES.get(theme, THEMES["corporate"])
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    total_slides = len(slides_data) + 2  # +title +closing

    # Title slide
    _build_title_slide(prs, title, subtitle, colors=colors)

    # Content slides
    for i, slide_data in enumerate(slides_data, 1):
        if slide_data.get("is_section"):
            _build_section_slide(prs, slide_data["title"], i, total_slides, colors=colors)
            continue

        # Generate image if requested
        image_bytes = None
        if with_images and slide_data.get("image_prompt"):
            prompt = (
                f"Professional presentation visual, clean modern style, "
                f"corporate design: {slide_data['image_prompt']}"
            )
            image_bytes = await generate_image(prompt, width=1280, height=960)

        # Check for special slide types
        has_quote = slide_data.get("quote")
        has_data = slide_data.get("data_highlights")

        if has_quote:
            _build_quote_slide(
                prs, slide_data["title"],
                quote_text=has_quote["text"],
                quote_author=has_quote["author"],
                bullets=slide_data.get("bullets", []),
                slide_num=i, total_slides=total_slides,
                image_bytes=image_bytes, colors=colors,
            )
        elif has_data:
            _build_data_slide(
                prs, slide_data["title"],
                data_highlights=has_data,
                bullets=slide_data.get("bullets", []),
                slide_num=i, total_slides=total_slides,
                image_bytes=image_bytes, colors=colors,
            )
        else:
            _build_content_slide(
                prs,
                title=slide_data["title"],
                bullets=slide_data.get("bullets", []),
                slide_num=i,
                total_slides=total_slides,
                image_bytes=image_bytes,
                colors=colors,
            )

    # Closing slide
    _build_closing_slide(prs, colors=colors)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def generate_presentation_sync(
    slides_data: list[dict[str, Any]],
    title: str,
    subtitle: str = "",
    with_images: bool = True,
    theme: str = "corporate",
) -> bytes:
    """Sync version of generate_presentation."""
    colors = THEMES.get(theme, THEMES["corporate"])
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    total_slides = len(slides_data) + 2

    _build_title_slide(prs, title, subtitle, colors=colors)

    for i, slide_data in enumerate(slides_data, 1):
        if slide_data.get("is_section"):
            _build_section_slide(prs, slide_data["title"], i, total_slides, colors=colors)
            continue

        image_bytes = None
        if with_images and slide_data.get("image_prompt"):
            prompt = (
                f"Professional presentation visual, clean modern style, "
                f"corporate design: {slide_data['image_prompt']}"
            )
            image_bytes = _generate_image_sync(prompt, width=1280, height=960)

        has_quote = slide_data.get("quote")
        has_data = slide_data.get("data_highlights")

        if has_quote:
            _build_quote_slide(
                prs, slide_data["title"],
                quote_text=has_quote["text"],
                quote_author=has_quote["author"],
                bullets=slide_data.get("bullets", []),
                slide_num=i, total_slides=total_slides,
                image_bytes=image_bytes, colors=colors,
            )
        elif has_data:
            _build_data_slide(
                prs, slide_data["title"],
                data_highlights=has_data,
                bullets=slide_data.get("bullets", []),
                slide_num=i, total_slides=total_slides,
                image_bytes=image_bytes, colors=colors,
            )
        else:
            _build_content_slide(
                prs,
                title=slide_data["title"],
                bullets=slide_data.get("bullets", []),
                slide_num=i,
                total_slides=total_slides,
                image_bytes=image_bytes,
                colors=colors,
            )

    _build_closing_slide(prs, colors=colors)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ── Prompt Builder for Agent ─────────────────────────────────────

def build_presentation_prompt(
    topic: str,
    slide_count: int = 10,
    language: str = "tr",
    research_context: str = "",
    mode: PresentationMode = "midi",
) -> str:
    """Build a rich prompt for the agent to generate structured slide content.
    Prompt complexity scales with mode."""
    lang_instruction = "Türkçe olarak yaz." if language == "tr" else "Write in English."

    research_block = ""
    if research_context:
        research_block = f"""

IMPORTANT — USE THIS RESEARCH DATA:
The following research has been gathered from real web sources. You MUST use this data
to create factual, data-rich slides. Include specific numbers, statistics, examples,
and cite sources where relevant.

{research_context}

"""

    # Mode-specific content instructions
    mode_instructions = {
        "mini": """
MINI MODE — Executive Summary:
- Focus on the TOP 3-5 most important points
- Each slide should deliver ONE key message
- Use impactful statistics and headlines
- Keep bullets to 3 per slide maximum
- Prioritize clarity over comprehensiveness
""",
        "midi": """
MIDI MODE — Standard Presentation:
- Balance depth with readability
- Include data, examples, and analysis
- Each slide should have 3-5 bullet points
- Add section dividers between major topics
- Include at least 2 slides with specific data/statistics
""",
        "maxi": """
MAXI MODE — Comprehensive Deep-Dive:
- Cover the topic exhaustively with multiple perspectives
- Include case studies, comparisons, and detailed analysis
- Each major section should have 3-5 slides
- Add data-heavy slides with statistics and market data
- Include pros/cons analysis, risk assessment, and recommendations
- Add a dedicated "Sources & References" slide at the end
- Use QUOTE: format for notable expert quotes
- Use DATA: format for key statistics to highlight
""",
    }

    extra_format_rules = ""
    if mode == "maxi":
        extra_format_rules = """
ADDITIONAL FORMAT OPTIONS (use where appropriate):
- QUOTE: "quote text" — Author Name (for notable quotes)
- DATA: metric_name | value | context (for key statistics to highlight visually)
"""

    return f"""Create a professional, data-rich presentation about: {topic}

Generate exactly {slide_count} content slides. {lang_instruction}

{mode_instructions.get(mode, mode_instructions["midi"])}

{research_block}

SLIDE DESIGN & CONTENT ARCHITECTURE:
The presentation engine will render your content with a professional geometric design system.
To maximize visual impact, follow these content architecture rules:

1. TITLE HIERARCHY: Each slide title should be punchy (max 8 words). Use "—" to add a subtitle hint.
   Good: "Yapay Zeka Pazarı — 2025 Büyüme Analizi"
   Bad: "Yapay Zeka Pazarının 2025 Yılındaki Büyüme Oranları ve Pazar Analizi Hakkında Bilgiler"

2. BULLET RHYTHM: Alternate between short punchy bullets and data-rich bullets.
   - Short: "Pazar büyümesi %37 (2025)"
   - Data-rich: "Kurumsal benimseme %72'ye ulaştı — en hızlı büyüyen segment üretken YZ"
   - Mix them: short, data, short, data, short

3. SECTION DIVIDERS: Use SECTION: between major topic shifts. These create bold visual breaks.
   Place them every 3-4 content slides for visual rhythm.

4. DATA SLIDES: When you have 2-3 key metrics, use DATA: format for visual card rendering.
   DATA: metric_name | value | brief context
   Example: DATA: Pazar Büyüklüğü | $190B | 2025 tahmini, %37 büyüme

5. QUOTE SLIDES: For expert opinions or notable statements, use QUOTE: format.
   QUOTE: "quote text" — Author Name
   These render as dramatic full-slide quotes with oversized quotation marks.

6. IMAGE PROMPTS: Be specific and visual. The engine generates real images from these.
   Bad: "AI image"
   Good: "Futuristic neural network visualization with glowing blue nodes connected by light beams, dark background, cinematic lighting"
   Include style cues: "infographic style", "isometric illustration", "cinematic", "minimal flat design"

7. VISUAL VARIETY: Don't make every slide the same type. Mix:
   - 60% content slides (bullets + image)
   - 15% section dividers
   - 15% data highlight slides (metrics cards)
   - 10% quote slides (if applicable)

FORMAT RULES (follow EXACTLY — this is CRITICAL for parsing):
- Start each slide with: SLIDE N: Title (e.g., SLIDE 1: Introduction)
- The word SLIDE and the number MUST appear at the start of the line
- Add bullet points with: - bullet text
- Add image description with: IMAGE: description for visual generation
- Use SECTION: Title for section divider slides (these don't count toward {slide_count})
- Do NOT use markdown headers (##, ###) for slides — use SLIDE N: format ONLY
- Do NOT wrap content in code blocks or markdown formatting
{extra_format_rules}

CONTENT QUALITY RULES:
- Each slide should have 3-5 bullet points maximum
- Bullets should be concise but informative (max 20 words each)
- Every slide MUST have an IMAGE line with a descriptive prompt for visual generation
- Image prompts should describe a specific, professional visual (not generic)
- INCLUDE real data: statistics, percentages, market sizes, growth rates from the research
- INCLUDE specific examples, company names, tool names, real-world applications
- INCLUDE year references (2024, 2025) for currency
- Add source references in bullets where relevant: (Kaynak: ...)
- Structure: Introduction → Key Concepts → Data & Statistics → Use Cases → Challenges → Future → Conclusion

EXAMPLE FORMAT:
SLIDE 1: Yapay Zeka Pazarı — Büyüklük ve Büyüme
- Global YZ pazarı 2025'te 190 milyar dolara ulaştı (%37 büyüme)
- 2030'a kadar 1.8 trilyon dolar bekleniyor (Kaynak: Statista)
- Kurumsal YZ benimseme oranı %72'ye yükseldi
- En hızlı büyüyen segment: Üretken YZ (%65 yıllık büyüme)
IMAGE: Global AI market growth chart with upward trend, modern infographic style, blue and green colors

SECTION: Temel Kavramlar

SLIDE 2: Çoklu Agent Sistemleri Nedir?
- Birden fazla YZ agent'ının koordineli çalıştığı mimari
- Her agent uzmanlaşmış bir göreve odaklanır (araştırma, analiz, kod)
- Orchestrator pattern: Merkezi yönetim ile görev dağılımı
- Gerçek dünya örneği: AutoGen, CrewAI, LangGraph
IMAGE: Multi-agent system architecture diagram with connected nodes, each representing a specialized AI agent

IMPORTANT: Start your response DIRECTLY with "SLIDE 1:" — do not add any introduction, explanation, or thinking before the slides.

Now generate the full presentation content using the research data provided:"""


# ── Topic Analysis for MINI/MIDI/MAXI ────────────────────────────

async def analyze_topic_for_presentation(
    topic: str,
    language: str = "tr",
) -> dict[str, Any]:
    """
    Perform quick research on a topic and return MINI/MIDI/MAXI options
    with descriptions of what each mode would cover.
    """
    from tools.search import web_search

    # Quick research to understand topic scope
    queries = [topic, f"{topic} overview 2025"]
    tasks = [web_search(q, max_results=5) for q in queries]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    # Count available data points
    fact_count = 0
    subtopics: set[str] = set()
    for result_set in results:
        if isinstance(result_set, Exception):
            continue
        for item in result_set:
            snippet = item.get("snippet", "")
            if snippet and len(snippet) > 30:
                fact_count += 1
            title = item.get("title", "")
            if title:
                subtopics.add(title[:60])

    # Build mode descriptions
    is_tr = language == "tr"
    options = {}
    for mode_key, cfg in MODE_CONFIG.items():
        lo, hi = cfg["slide_range"]
        desc = cfg["description_tr"] if is_tr else cfg["description_en"]
        options[mode_key] = {
            "label": cfg["label"],
            "emoji": cfg["emoji"],
            "slide_range": f"{lo}-{hi}",
            "default_slides": cfg["default_slides"],
            "description": desc,
            "research_queries": cfg["research_queries"],
        }

    return {
        "topic": topic,
        "available_data_points": fact_count,
        "discovered_subtopics": list(subtopics)[:10],
        "options": options,
        "language": language,
    }


def format_analysis_response(analysis: dict[str, Any]) -> str:
    """Format the topic analysis into a user-friendly message with MINI/MIDI/MAXI options."""
    topic = analysis["topic"]
    data_points = analysis["available_data_points"]
    subtopics = analysis.get("discovered_subtopics", [])
    options = analysis["options"]

    lines = [
        f"🔍 **Konu Analizi:** {topic}",
        f"📊 **Bulunan Veri Noktası:** {data_points}",
    ]

    if subtopics:
        lines.append(f"📌 **Keşfedilen Alt Konular:** {', '.join(subtopics[:5])}")

    lines.append("")
    lines.append("🎯 **Sunum Modu Seçenekleri:**")
    lines.append("")

    for mode_key in ("mini", "midi", "maxi"):
        opt = options[mode_key]
        lines.append(
            f"{opt['emoji']} **{opt['label']}** ({opt['slide_range']} slayt) — "
            f"{opt['description']}"
        )
        lines.append(
            f"   🔬 {opt['research_queries']} araştırma sorgusu kullanılacak"
        )
        lines.append("")

    lines.append("💬 **Seçimini yaz:** `MINI`, `MIDI` veya `MAXI`")
    lines.append("   (Veya doğrudan slayt sayısı belirtebilirsin: örn. `MIDI 15`)")

    return "\n".join(lines)
