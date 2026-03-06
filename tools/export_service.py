"""
Export Service — Professional PDF & Markdown export with Turkish character support.
Uses ReportLab for PDF generation with proper pagination, headers, and Unicode fonts.
"""

from __future__ import annotations

import io
import os
import re
from pathlib import Path
from typing import Any

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm, mm
from reportlab.lib.colors import HexColor
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak,
    Table, TableStyle, HRFlowable,
)
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont


# ── Font Registration (Turkish character support) ────────────────

_FONT_REGISTERED = False

def _register_fonts():
    """Register a Unicode-capable font for Turkish characters.

    Search order: reportlab bundled Vera (always available) → Windows system
    fonts → Linux system fonts → Helvetica fallback (limited Turkish support).
    """
    global _FONT_REGISTERED
    if _FONT_REGISTERED:
        return

    import reportlab as _rl
    _rl_fonts_dir = Path(_rl.__file__).parent / "fonts"

    font_pairs = [
        # 1. reportlab bundled Vera — always available, full Unicode/Turkish
        (str(_rl_fonts_dir / "Vera.ttf"), str(_rl_fonts_dir / "VeraBd.ttf")),
        # 2. Project /fonts directory (optional bundled DejaVu)
        (str(Path(__file__).parent.parent / "fonts" / "DejaVuSans.ttf"),
         str(Path(__file__).parent.parent / "fonts" / "DejaVuSans-Bold.ttf")),
        # 3. Windows system fonts
        ("C:/Windows/Fonts/arial.ttf", "C:/Windows/Fonts/arialbd.ttf"),
        ("C:/Windows/Fonts/calibri.ttf", "C:/Windows/Fonts/calibrib.ttf"),
        ("C:/Windows/Fonts/segoeui.ttf", "C:/Windows/Fonts/segoeuib.ttf"),
        ("C:/Windows/Fonts/tahoma.ttf", "C:/Windows/Fonts/tahomabd.ttf"),
        # 4. Linux / Docker system fonts
        ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
         "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        ("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
         "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"),
        ("/usr/share/fonts/truetype/freefont/FreeSans.ttf",
         "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf"),
        # 5. macOS
        ("/Library/Fonts/Arial.ttf", "/Library/Fonts/Arial Bold.ttf"),
    ]

    for regular, bold in font_pairs:
        if not os.path.exists(regular):
            continue
        try:
            pdfmetrics.registerFont(TTFont("TRFont", regular))
            if os.path.exists(bold):
                pdfmetrics.registerFont(TTFont("TRFontBold", bold))
            else:
                pdfmetrics.registerFont(TTFont("TRFontBold", regular))
            _FONT_REGISTERED = True
            return
        except Exception:
            continue

    # Final fallback: reportlab's built-in Helvetica (no Turkish ş/ğ/ı/ç/ö/ü)
    _FONT_REGISTERED = True


def _get_font_name(bold: bool = False) -> str:
    """Get registered font name."""
    _register_fonts()
    if bold:
        try:
            pdfmetrics.getFont("TRFontBold")
            return "TRFontBold"
        except KeyError:
            pass
    try:
        pdfmetrics.getFont("TRFont")
        return "TRFont"
    except KeyError:
        return "Helvetica-Bold" if bold else "Helvetica"


# ── PDF Styles ───────────────────────────────────────────────────

def _build_styles() -> dict[str, ParagraphStyle]:
    """Build PDF paragraph styles with Turkish font support."""
    font = _get_font_name()
    font_bold = _get_font_name(bold=True)

    return {
        "title": ParagraphStyle(
            "CustomTitle",
            fontName=font_bold,
            fontSize=22,
            leading=28,
            spaceAfter=12,
            textColor=HexColor("#1a1a2e"),
            alignment=TA_CENTER,
        ),
        "subtitle": ParagraphStyle(
            "CustomSubtitle",
            fontName=font,
            fontSize=11,
            leading=14,
            spaceAfter=20,
            textColor=HexColor("#666666"),
            alignment=TA_CENTER,
        ),
        "h1": ParagraphStyle(
            "CustomH1",
            fontName=font_bold,
            fontSize=16,
            leading=22,
            spaceBefore=18,
            spaceAfter=8,
            textColor=HexColor("#16213e"),
        ),
        "h2": ParagraphStyle(
            "CustomH2",
            fontName=font_bold,
            fontSize=13,
            leading=18,
            spaceBefore=14,
            spaceAfter=6,
            textColor=HexColor("#0f3460"),
        ),
        "h3": ParagraphStyle(
            "CustomH3",
            fontName=font_bold,
            fontSize=11,
            leading=15,
            spaceBefore=10,
            spaceAfter=4,
            textColor=HexColor("#533483"),
        ),
        "body": ParagraphStyle(
            "CustomBody",
            fontName=font,
            fontSize=10,
            leading=14,
            spaceAfter=6,
            textColor=HexColor("#333333"),
            alignment=TA_JUSTIFY,
        ),
        "bullet": ParagraphStyle(
            "CustomBullet",
            fontName=font,
            fontSize=10,
            leading=14,
            spaceAfter=3,
            leftIndent=20,
            textColor=HexColor("#333333"),
        ),
        "code": ParagraphStyle(
            "CustomCode",
            fontName="Courier",
            fontSize=8,
            leading=11,
            spaceAfter=6,
            leftIndent=10,
            backColor=HexColor("#f5f5f5"),
            textColor=HexColor("#2d2d2d"),
        ),
        "footer": ParagraphStyle(
            "CustomFooter",
            fontName=font,
            fontSize=8,
            textColor=HexColor("#999999"),
            alignment=TA_CENTER,
        ),
    }


# ── Markdown to PDF Elements ────────────────────────────────────

def _escape_xml(text: str) -> str:
    """Escape XML special characters for ReportLab paragraphs."""
    text = text.replace("&", "&amp;")
    text = text.replace("<", "&lt;")
    text = text.replace(">", "&gt;")
    return text


def _md_to_flowables(markdown_text: str, styles: dict) -> list:
    """Convert markdown text to ReportLab flowable elements."""
    flowables = []
    lines = markdown_text.split("\n")
    in_code_block = False
    code_buffer = []

    for line in lines:
        stripped = line.strip()

        # Code block toggle
        if stripped.startswith("```"):
            if in_code_block:
                # End code block
                code_text = _escape_xml("\n".join(code_buffer))
                if code_text.strip():
                    flowables.append(Paragraph(code_text.replace("\n", "<br/>"), styles["code"]))
                    flowables.append(Spacer(1, 4))
                code_buffer = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_buffer.append(line)
            continue

        # Empty line
        if not stripped:
            flowables.append(Spacer(1, 4))
            continue

        # Horizontal rule
        if stripped in ("---", "***", "___"):
            flowables.append(HRFlowable(width="100%", thickness=0.5, color=HexColor("#cccccc")))
            flowables.append(Spacer(1, 6))
            continue

        # Headers
        if stripped.startswith("### "):
            text = _escape_xml(stripped[4:])
            flowables.append(Paragraph(text, styles["h3"]))
            continue
        if stripped.startswith("## "):
            text = _escape_xml(stripped[3:])
            flowables.append(Paragraph(text, styles["h2"]))
            continue
        if stripped.startswith("# "):
            text = _escape_xml(stripped[2:])
            flowables.append(Paragraph(text, styles["h1"]))
            continue

        # Bullet points
        if stripped.startswith(("- ", "* ", "• ")):
            text = _escape_xml(stripped[2:])
            flowables.append(Paragraph(f"• {text}", styles["bullet"]))
            continue

        # Numbered lists
        num_match = re.match(r"^(\d+)\.\s+(.+)", stripped)
        if num_match:
            num, text = num_match.groups()
            text = _escape_xml(text)
            flowables.append(Paragraph(f"{num}. {text}", styles["bullet"]))
            continue

        # Checkbox items
        if stripped.startswith(("- [ ] ", "- [x] ", "- [X] ")):
            checked = stripped[3] in ("x", "X")
            mark = "☑" if checked else "☐"
            text = _escape_xml(stripped[6:])
            flowables.append(Paragraph(f"{mark} {text}", styles["bullet"]))
            continue

        # Bold text inline
        text = _escape_xml(stripped)
        # Convert **bold** to <b>bold</b>
        text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", text)

        flowables.append(Paragraph(text, styles["body"]))

    return flowables


# ── PDF Generation ───────────────────────────────────────────────

def _header_footer(canvas, doc):
    """Draw header and footer on each page."""
    canvas.saveState()
    font = _get_font_name()

    # Footer: page number
    canvas.setFont(font, 8)
    canvas.setFillColor(HexColor("#999999"))
    page_num = f"Sayfa {doc.page}"
    canvas.drawCentredString(A4[0] / 2, 15 * mm, page_num)

    # Header line (not on first page)
    if doc.page > 1:
        canvas.setStrokeColor(HexColor("#e0e0e0"))
        canvas.setLineWidth(0.5)
        canvas.line(2 * cm, A4[1] - 1.5 * cm, A4[0] - 2 * cm, A4[1] - 1.5 * cm)

    canvas.restoreState()


def generate_pdf(markdown_content: str, title: str = "Rapor") -> bytes:
    """
    Generate a professional PDF from markdown content.
    Returns PDF as bytes.
    """
    _register_fonts()
    styles = _build_styles()

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        topMargin=2 * cm,
        bottomMargin=2.5 * cm,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        title=title,
        author="Multi-Agent Ops Center",
    )

    flowables = []

    # Title page elements
    flowables.append(Spacer(1, 3 * cm))
    flowables.append(Paragraph(_escape_xml(title), styles["title"]))
    flowables.append(Spacer(1, 0.5 * cm))

    from datetime import datetime
    date_str = datetime.now().strftime("%d %B %Y, %H:%M")
    flowables.append(Paragraph(f"Oluşturulma: {date_str}", styles["subtitle"]))
    flowables.append(Paragraph("Multi-Agent Ops Center", styles["subtitle"]))
    flowables.append(Spacer(1, 1 * cm))
    flowables.append(HRFlowable(width="60%", thickness=1, color=HexColor("#0f3460")))
    flowables.append(Spacer(1, 1 * cm))

    # Content
    content_flowables = _md_to_flowables(markdown_content, styles)
    flowables.extend(content_flowables)

    doc.build(flowables, onFirstPage=_header_footer, onLaterPages=_header_footer)
    return buffer.getvalue()

def generate_html(markdown_content: str, title: str = "Rapor") -> str:
    """
    Convert markdown content to a standalone, styled HTML report.
    Supports Turkish characters, images, tables, and code blocks.
    Returns HTML string (UTF-8).
    """
    import html as html_lib
    from datetime import datetime

    date_str = datetime.now().strftime("%d %B %Y, %H:%M")

    def md_to_html(text: str) -> str:
        lines = text.split("\n")
        out = []
        in_code = False
        code_lang = ""
        code_buf: list[str] = []
        in_table = False
        table_buf: list[str] = []

        def flush_table():
            nonlocal in_table, table_buf
            if not table_buf:
                return
            rows = [r for r in table_buf if r.strip()]
            html_rows = []
            for i, row in enumerate(rows):
                cells = [c.strip() for c in row.strip("|").split("|")]
                if i == 1 and all(set(c.replace("-", "").replace(":", "").replace(" ", "")) == set() for c in cells):
                    continue  # separator row
                tag = "th" if i == 0 else "td"
                html_rows.append("<tr>" + "".join(f"<{tag}>{html_lib.escape(c)}</{tag}>" for c in cells) + "</tr>")
            out.append('<div class="table-wrap"><table>' + "".join(html_rows) + "</table></div>")
            table_buf = []
            in_table = False

        for line in lines:
            stripped = line.strip()

            # Code block
            if stripped.startswith("```"):
                if in_code:
                    code_text = html_lib.escape("\n".join(code_buf))
                    out.append(f'<pre><code class="language-{html_lib.escape(code_lang)}">{code_text}</code></pre>')
                    code_buf = []
                    in_code = False
                    code_lang = ""
                else:
                    if in_table:
                        flush_table()
                    in_code = True
                    code_lang = stripped[3:].strip()
                continue

            if in_code:
                code_buf.append(line)
                continue

            # Table
            if stripped.startswith("|"):
                in_table = True
                table_buf.append(stripped)
                continue
            elif in_table:
                flush_table()

            if not stripped:
                out.append("<br>")
                continue

            if stripped in ("---", "***", "___"):
                out.append("<hr>")
                continue

            # Image: ![alt](url)
            img_match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", stripped)
            if img_match:
                alt, src = img_match.groups()
                out.append(f'<figure><img src="{html_lib.escape(src)}" alt="{html_lib.escape(alt)}" loading="lazy"><figcaption>{html_lib.escape(alt)}</figcaption></figure>')
                continue

            def inline(t: str) -> str:
                # Bold+italic
                t = re.sub(r"\*\*\*(.+?)\*\*\*", r"<strong><em>\1</em></strong>", t)
                # Bold
                t = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", t)
                t = re.sub(r"__(.+?)__", r"<strong>\1</strong>", t)
                # Italic
                t = re.sub(r"\*(.+?)\*", r"<em>\1</em>", t)
                t = re.sub(r"_(.+?)_", r"<em>\1</em>", t)
                # Inline code
                t = re.sub(r"`([^`]+)`", lambda m: f"<code>{html_lib.escape(m.group(1))}</code>", t)
                # Links
                t = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r'<a href="\2" target="_blank">\1</a>', t)
                return t

            esc = html_lib.escape(stripped)

            if stripped.startswith("#### "):
                out.append(f"<h4>{inline(html_lib.escape(stripped[5:]))}</h4>")
            elif stripped.startswith("### "):
                out.append(f"<h3>{inline(html_lib.escape(stripped[4:]))}</h3>")
            elif stripped.startswith("## "):
                out.append(f"<h2>{inline(html_lib.escape(stripped[3:]))}</h2>")
            elif stripped.startswith("# "):
                out.append(f"<h1>{inline(html_lib.escape(stripped[2:]))}</h1>")
            elif stripped.startswith(("- ", "* ", "• ")):
                out.append(f"<li>{inline(html_lib.escape(stripped[2:]))}</li>")
            elif re.match(r"^\d+\.\s", stripped):
                m = re.match(r"^(\d+)\.\s+(.+)", stripped)
                if m:
                    out.append(f"<li>{inline(html_lib.escape(m.group(2)))}</li>")
            elif stripped.startswith("> "):
                out.append(f"<blockquote>{inline(html_lib.escape(stripped[2:]))}</blockquote>")
            else:
                out.append(f"<p>{inline(esc)}</p>")

        if in_table:
            flush_table()
        if in_code and code_buf:
            out.append(f'<pre><code>{html_lib.escape(chr(10).join(code_buf))}</code></pre>')

        return "\n".join(out)

    body = md_to_html(markdown_content)
    safe_title = html_lib.escape(title)

    return f"""<!DOCTYPE html>
<html lang="tr">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{safe_title}</title>
<style>
  :root {{
    --primary: #0f3460;
    --accent: #533483;
    --bg: #f8f9fa;
    --surface: #ffffff;
    --text: #1a1a2e;
    --muted: #666;
    --border: #e0e0e0;
    --code-bg: #f5f5f5;
  }}
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ font-family: 'Segoe UI', Arial, sans-serif; background: var(--bg); color: var(--text); line-height: 1.7; }}
  .container {{ max-width: 860px; margin: 0 auto; padding: 2rem 1.5rem; }}
  header {{ background: var(--primary); color: #fff; padding: 2.5rem 2rem; border-radius: 12px; margin-bottom: 2rem; }}
  header h1 {{ font-size: 1.8rem; font-weight: 700; margin-bottom: 0.4rem; }}
  header .meta {{ font-size: 0.85rem; opacity: 0.75; }}
  .content {{ background: var(--surface); border-radius: 12px; padding: 2rem; box-shadow: 0 2px 12px rgba(0,0,0,0.06); }}
  h1 {{ font-size: 1.6rem; color: var(--primary); margin: 1.5rem 0 0.75rem; border-bottom: 2px solid var(--border); padding-bottom: 0.4rem; }}
  h2 {{ font-size: 1.3rem; color: var(--primary); margin: 1.4rem 0 0.6rem; }}
  h3 {{ font-size: 1.1rem; color: var(--accent); margin: 1.2rem 0 0.5rem; }}
  h4 {{ font-size: 1rem; color: var(--muted); margin: 1rem 0 0.4rem; }}
  p {{ margin: 0.6rem 0; }}
  li {{ margin: 0.3rem 0 0.3rem 1.5rem; list-style: disc; }}
  ol li {{ list-style: decimal; }}
  blockquote {{ border-left: 4px solid var(--accent); padding: 0.5rem 1rem; margin: 1rem 0; background: #f3f0ff; border-radius: 0 8px 8px 0; color: var(--accent); font-style: italic; }}
  code {{ background: var(--code-bg); padding: 0.15em 0.4em; border-radius: 4px; font-family: 'Consolas', monospace; font-size: 0.88em; color: #c7254e; }}
  pre {{ background: #1e1e2e; color: #cdd6f4; padding: 1.2rem; border-radius: 8px; overflow-x: auto; margin: 1rem 0; }}
  pre code {{ background: none; color: inherit; padding: 0; font-size: 0.85rem; }}
  .table-wrap {{ overflow-x: auto; margin: 1rem 0; }}
  table {{ border-collapse: collapse; width: 100%; font-size: 0.9rem; }}
  th, td {{ border: 1px solid var(--border); padding: 0.6rem 0.9rem; text-align: left; }}
  th {{ background: var(--primary); color: #fff; font-weight: 600; }}
  tr:nth-child(even) td {{ background: #f5f7ff; }}
  hr {{ border: none; border-top: 1px solid var(--border); margin: 1.5rem 0; }}
  figure {{ margin: 1.5rem 0; text-align: center; }}
  figure img {{ max-width: 100%; border-radius: 8px; box-shadow: 0 2px 12px rgba(0,0,0,0.12); }}
  figcaption {{ font-size: 0.82rem; color: var(--muted); margin-top: 0.4rem; }}
  a {{ color: var(--accent); text-decoration: none; }}
  a:hover {{ text-decoration: underline; }}
  footer {{ text-align: center; color: var(--muted); font-size: 0.8rem; margin-top: 2rem; padding-top: 1rem; border-top: 1px solid var(--border); }}
  @media print {{
    body {{ background: #fff; }}
    .container {{ padding: 0; }}
    header {{ border-radius: 0; }}
    .content {{ box-shadow: none; }}
  }}
</style>
</head>
<body>
<div class="container">
  <header>
    <h1>{safe_title}</h1>
    <div class="meta">Oluşturulma: {date_str} &nbsp;|&nbsp; Multi-Agent Ops Center</div>
  </header>
  <div class="content">
{body}
  </div>
  <footer>Multi-Agent Ops Center &copy; {datetime.now().year}</footer>
</div>
</body>
</html>"""


def generate_presentation_pdf(pptx_path: str, title: str = "Sunum") -> bytes:
    """
    Read a PPTX file and generate a professional PDF with slide content.
    Each slide becomes a section in the PDF with title + bullets.
    """
    from pptx import Presentation as PptxPresentation

    _register_fonts()
    styles = _build_styles()

    font = _get_font_name()
    font_bold = _get_font_name(bold=True)
    slide_title_style = ParagraphStyle(
        "SlideTitle",
        fontName=font_bold,
        fontSize=14,
        leading=20,
        spaceBefore=16,
        spaceAfter=8,
        textColor=HexColor("#0f3460"),
    )
    slide_num_style = ParagraphStyle(
        "SlideNum",
        fontName=font_bold,
        fontSize=9,
        leading=12,
        textColor=HexColor("#999999"),
    )
    slide_body_style = ParagraphStyle(
        "SlideBody",
        fontName=font,
        fontSize=10,
        leading=15,
        spaceAfter=4,
        leftIndent=15,
        textColor=HexColor("#333333"),
    )

    prs = PptxPresentation(pptx_path)

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        topMargin=2 * cm,
        bottomMargin=2.5 * cm,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        title=title,
        author="Multi-Agent Ops Center",
    )

    flowables = []

    # Title page
    flowables.append(Spacer(1, 3 * cm))
    flowables.append(Paragraph(_escape_xml(title), styles["title"]))
    flowables.append(Spacer(1, 0.5 * cm))

    from datetime import datetime
    date_str = datetime.now().strftime("%d %B %Y, %H:%M")
    flowables.append(Paragraph(f"Oluşturulma: {date_str}", styles["subtitle"]))
    flowables.append(Paragraph(f"{len(prs.slides)} Slayt | AI Destekli Sunum", styles["subtitle"]))
    flowables.append(Spacer(1, 1 * cm))
    flowables.append(HRFlowable(width="60%", thickness=1, color=HexColor("#0f3460")))
    flowables.append(PageBreak())

    # Extract slides
    for i, slide in enumerate(prs.slides, 1):
        flowables.append(Paragraph(f"Slayt {i}/{len(prs.slides)}", slide_num_style))

        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is None:
                continue
            for para in text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)

        if texts:
            slide_title_text = texts[0]
            body_texts = texts[1:]
        else:
            slide_title_text = f"Slayt {i}"
            body_texts = []

        flowables.append(Paragraph(_escape_xml(slide_title_text), slide_title_style))

        for text in body_texts:
            flowables.append(Paragraph(f"• {_escape_xml(text)}", slide_body_style))

        flowables.append(Spacer(1, 6))
        flowables.append(HRFlowable(width="100%", thickness=0.3, color=HexColor("#e0e0e0")))
        flowables.append(Spacer(1, 4))

    doc.build(flowables, onFirstPage=_header_footer, onLaterPages=_header_footer)
    return buffer.getvalue()
