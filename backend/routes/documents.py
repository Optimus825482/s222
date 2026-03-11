"""
Doküman yükleme, metin çıkarma ve dosya yönetimi endpointleri.
Desteklenen formatlar: PDF, DOCX, XLSX, PPTX, PNG, JPG, JPEG, GIF, WEBP
"""

import logging
import sys
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from pydantic import BaseModel

_parent = str(Path(__file__).parent.parent)
if _parent not in sys.path:
    sys.path.insert(0, _parent)

from deps import get_current_user, _audit
from config import DATA_DIR

logger = logging.getLogger(__name__)

router = APIRouter(tags=["documents"])

# ── Sabitler ─────────────────────────────────────────────────────

UPLOADS_DIR = DATA_DIR / "uploads"
UPLOADS_DIR.mkdir(parents=True, exist_ok=True)

MAX_FILE_SIZE = 20 * 1024 * 1024  # 20 MB
MAX_EXTRACTED_CHARS = 100_000

ALLOWED_EXTENSIONS: set[str] = {
    ".pdf", ".docx", ".xlsx", ".pptx",
    ".png", ".jpg", ".jpeg", ".gif", ".webp",
}

IMAGE_EXTENSIONS: set[str] = {".png", ".jpg", ".jpeg", ".gif", ".webp"}

EXTENSION_TO_CONTENT_TYPE: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}


# ── Pydantic Modelleri ───────────────────────────────────────────

class ExtractRequest(BaseModel):
    file_id: str


# ── Metin Çıkarma Fonksiyonları ─────────────────────────────────

def _extract_pdf(file_path: Path) -> str:
    """PDF dosyasından tüm sayfaların metnini çıkar."""
    from pypdf import PdfReader

    reader = PdfReader(str(file_path))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_docx(file_path: Path) -> str:
    """DOCX dosyasından paragraf ve tablo hücrelerini çıkar."""
    from docx import Document

    doc = Document(str(file_path))
    parts: list[str] = []

    # Paragraflar
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # Tablo hücreleri
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    return "\n".join(parts)


def _extract_xlsx(file_path: Path) -> str:
    """XLSX dosyasından tüm sayfaları ve satırları çıkar."""
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells).strip()
            if line.replace("|", "").strip():
                parts.append(line)

    wb.close()
    return "\n".join(parts)


def _extract_pptx(file_path: Path) -> str:
    """PPTX dosyasından tüm slaytlardaki metin kutularını çıkar."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            text_frame = getattr(shape, "text_frame", None)
            if getattr(shape, "has_text_frame", False) and text_frame is not None:
                for paragraph in text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"--- Slayt {slide_idx} ---")
            parts.extend(slide_texts)

    return "\n".join(parts)


def _extract_text(file_path: Path, extension: str) -> str:
    """Dosya türüne göre metin çıkarma — güvenli, truncated."""
    if extension in IMAGE_EXTENSIONS:
        return "[Image file — visual analysis required]"

    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
    }

    extractor = extractors.get(extension)
    if not extractor:
        return f"[Unsupported format: {extension}]"

    try:
        text = extractor(file_path)
        # Güvenlik: max karakter sınırı
        if len(text) > MAX_EXTRACTED_CHARS:
            text = text[:MAX_EXTRACTED_CHARS] + f"\n\n[Truncated — {MAX_EXTRACTED_CHARS:,} karakter sınırına ulaşıldı]"
        return text
    except Exception as e:
        logger.error(f"Metin çıkarma hatası ({file_path.name}): {e}")
        return f"[Extraction failed: {type(e).__name__}]"


# ── Yardımcı Fonksiyonlar ───────────────────────────────────────

def _validate_extension(filename: str) -> str:
    """Dosya uzantısını doğrula, güvenli extension döndür."""
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Desteklenmeyen dosya formatı: {ext}. İzin verilen: {', '.join(sorted(ALLOWED_EXTENSIONS))}",
        )
    return ext


def _build_file_metadata(file_id: str, original_name: str, ext: str, size: int, saved_path: Path) -> dict[str, Any]:
    """Upload metadata sözlüğü oluştur."""
    return {
        "file_id": file_id,
        "filename": original_name,
        "content_type": EXTENSION_TO_CONTENT_TYPE.get(ext, "application/octet-stream"),
        "size_bytes": size,
        "extension": ext,
        "saved_path": str(saved_path),
        "uploaded_at": datetime.now(timezone.utc).isoformat(),
    }


# In-memory dosya kayıt defteri (basit, production'da DB'ye taşınır)
_FILE_REGISTRY: dict[str, dict[str, Any]] = {}


# ── Endpointler ──────────────────────────────────────────────────

@router.post("/api/documents/upload")
async def upload_document(
    file: UploadFile = File(...),
    user: dict = Depends(get_current_user),
):
    """Doküman yükle ve otomatik metin çıkar.

    Desteklenen formatlar: PDF, DOCX, XLSX, PPTX, PNG, JPG, JPEG, GIF, WEBP
    Maksimum dosya boyutu: 20 MB
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="Dosya adı gerekli")

    ext = _validate_extension(file.filename)

    # Dosya boyutu kontrolü — tamamını belleğe oku
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"Dosya boyutu çok büyük. Maksimum: {MAX_FILE_SIZE // (1024 * 1024)} MB",
        )
    if len(content) == 0:
        raise HTTPException(status_code=400, detail="Boş dosya yüklenemez")

    # UUID tabanlı dosya adı — path traversal koruması
    file_id = uuid.uuid4().hex
    safe_filename = f"{file_id}{ext}"
    saved_path = UPLOADS_DIR / safe_filename

    # Diske yaz
    saved_path.write_bytes(content)

    # Metin çıkar
    extracted_text = _extract_text(saved_path, ext)

    # Kayıt defterine ekle
    metadata = _build_file_metadata(file_id, file.filename, ext, len(content), saved_path)
    _FILE_REGISTRY[file_id] = metadata

    _audit("document_upload", user["user_id"], detail=f"{file.filename} ({len(content)} bytes)")
    logger.info(f"Doküman yüklendi: {file.filename} → {safe_filename} ({len(content)} bytes)")

    return {
        "file_id": file_id,
        "filename": file.filename,
        "content_type": metadata["content_type"],
        "size_bytes": len(content),
        "extracted_text": extracted_text,
        "char_count": len(extracted_text),
    }


@router.post("/api/documents/extract")
async def extract_document_text(
    body: ExtractRequest,
    user: dict = Depends(get_current_user),
):
    """Daha önce yüklenmiş bir dosyadan metin çıkar."""
    meta = _FILE_REGISTRY.get(body.file_id)
    if not meta:
        raise HTTPException(status_code=404, detail="Dosya bulunamadı")

    saved_path = Path(meta["saved_path"])
    if not saved_path.exists():
        raise HTTPException(status_code=404, detail="Dosya diskte bulunamadı")

    extracted_text = _extract_text(saved_path, meta["extension"])

    _audit("document_extract", user["user_id"], detail=f"file_id={body.file_id}")

    return {
        "file_id": body.file_id,
        "extracted_text": extracted_text,
        "char_count": len(extracted_text),
    }


@router.get("/api/documents/uploads")
async def list_uploaded_documents(user: dict = Depends(get_current_user)):
    """Yüklenmiş dosyaları listele."""
    _audit("document_list", user["user_id"])

    files = [
        {
            "file_id": meta["file_id"],
            "filename": meta["filename"],
            "content_type": meta["content_type"],
            "size_bytes": meta["size_bytes"],
            "uploaded_at": meta["uploaded_at"],
        }
        for meta in _FILE_REGISTRY.values()
    ]

    # En yeni dosya en üstte
    files.sort(key=lambda f: f["uploaded_at"], reverse=True)

    return {"files": files}


@router.delete("/api/documents/uploads/{file_id}")
async def delete_uploaded_document(
    file_id: str,
    user: dict = Depends(get_current_user),
):
    """Yüklenmiş dosyayı sil (disk + kayıt)."""
    meta = _FILE_REGISTRY.pop(file_id, None)
    if not meta:
        raise HTTPException(status_code=404, detail="Dosya bulunamadı")

    # Diskten sil
    saved_path = Path(meta["saved_path"])
    try:
        if saved_path.exists():
            saved_path.unlink()
    except OSError as e:
        logger.warning(f"Dosya silinemedi ({saved_path}): {e}")

    _audit("document_delete", user["user_id"], detail=f"file_id={file_id}, name={meta['filename']}")

    return {"deleted": True, "file_id": file_id, "filename": meta["filename"]}
